from __future__ import annotations

import io
import tempfile
import urllib.error
import urllib.request
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "monkey_collector_collection_pipeline.pptx"
KROKI_URL = "https://kroki.io/mermaid/png"
# kroki fronts Cloudflare bot protection; the default urllib UA gets a 403.
# A browser UA is required for the POST render endpoint to return a PNG.
USER_AGENT = (
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
    "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0 Safari/537.36"
)
SLIDE_W = 13.333
SLIDE_H = 7.5


@dataclass(frozen=True)
class SlideSpec:
    title: str
    bullets: list[str]
    notes: list[str]
    mermaid: str | None = None


DIAGRAM_A = """\
flowchart TD
  EV[AccessibilityEvent] --> ST[ScreenStabilizer]
  ST -->|no-change| N[N signal]
  ST -->|external app| E[E signal]
  ST -->|visual change| CAP[screenshot + XML + metadata over TCP]
  CAP --> SRV[Python Server: consume latest signal, parse XML]
  SRV --> SM[ScreenMatcher: BM25 unique-page matching LLM-free]
  SM -->|pending: no interactable| SKIP[no page node, no observation]
  SM -->|page_key| EXP[LLMGuidedExplorer: select action]
  EXP --> ADB[ADB execute action]
  ADB --> LOG[events.jsonl: frame_index / page_key / observation_num]
  LOG --> EV
"""

DIAGRAM_B = """\
flowchart TD
  IN[encoded XML + screenshot] --> PF[structure fingerprint pre-filter]
  PF -->|exact revisit| REUSE[reuse cached page_key + observation_num]
  PF -->|0 interactable button/input| PEND[PENDING: reject, no node]
  PF -->|new or uncached| SER[serialize element-line document]
  SER --> BM[BM25 top-K candidate pages]
  BM --> VER{element diff < max AND pixel gate < 0.3}
  VER -->|pass| MERGE[BM25_MERGE: stored page_key]
  VER -->|all fail| NEW[NEW: new page_key, add to BM25 corpus]
  MERGE --> OBS[observation identify: page-local luminance dedup]
  NEW --> OBS0[observation 0]
"""

DIAGRAM_C = """\
flowchart TD
  S[select_action per step] --> Q1{navigation in progress?}
  Q1 -->|yes| NAV[re-match next queued step by signature and execute]
  Q1 -->|no| Q2{current screen has unexplored action?}
  Q2 -->|yes| CUR[pick one unexplored, long_touch last]
  Q2 -->|no| Q3{global unexplored target exists?}
  Q3 -->|yes| PLAN[strategy BFS/DFS/GREEDY picks target, TransitionGraph shortest path, execute first step]
  Q3 -->|no| BACK[PressBack; root screen: in-screen tap instead]
"""


SLIDES = [
    SlideSpec(
        title="Monkey-Collector 수집 방식",
        bullets=[
            "현재 코드·세팅 기준 파이프라인 요약",
            "근거 문서: ARCHITECTURE.md 중심, README.md·config/run.yaml 보조",
            "수집은 App 전환 감지와 Python Server 액션 선택·저장으로 분리된다",
        ],
        notes=[
            "ARCHITECTURE.md 전체를 요약한 발표 표지다.",
            "다음 슬라이드부터 각 섹션 근거를 순서대로 설명한다.",
        ],
    ),
    SlideSpec(
        title="시스템 개요와 역할 분리",
        bullets=[
            "Android App은 foreground 감지, 안정화 판단, screenshot·XML 전송을 맡는다.",
            "Python Server는 TCP 수신, XML 파싱, action 선택, ADB 실행, 세션 저장을 맡는다.",
            "전환 감지는 App, no-change·first screen·external recovery는 server loop가 담당한다.",
            "세션 저장은 data/{package}와 runtime/{package} 두 root로 분리된다.",
        ],
        notes=[
            "근거: §1 시스템 개요, §2 컴포넌트 구조.",
            "핵심 메시지는 App과 Server의 책임 경계가 명확하다는 점이다.",
        ],
        mermaid="""\
flowchart LR
  APP[Android App\nAccessibilityService] -->|TCP signals + payload| TCP[TCP channel]
  TCP --> SRV[Python Server]
  APP --> APP2[Detect foreground\nStabilize screen\nCapture screenshot/XML]
  SRV --> SRV2[Parse XML\nSelect action\nExecute via ADB\nPersist session]
""",
    ),
    SlideSpec(
        title="Server-driven standby 루프와 세션 핸드셰이크",
        bullets=[
            "클라이언트는 standby 연결을 유지하다가 START 메시지를 받으면 수집을 시작한다.",
            "서버는 앱마다 force-stop과 launch 뒤 collection loop를 실행한다.",
            "세션 종료 시 SESSION_END 이후 클라이언트가 F를 보내고 fresh socket으로 재접속한다.",
            "reset_for_new_session은 큐·이벤트만 초기화하고 소켓은 유지해야 한다.",
        ],
        notes=[
            "근거: §2 Android App, §3 세션 전환 핸드셰이크.",
            "두 번째 세션부터 타임아웃이 나지 않게 소켓을 닫지 않는 계약이 중요하다.",
        ],
        mermaid="""\
sequenceDiagram
  participant S as Python Server
  participant C as Android Client
  S->>C: {"type":"START","package":"pkg"}
  C->>S: P ACK
  S->>S: force_stop(pkg) + launch_app(pkg)
  S->>S: run_collection_loop
  S->>C: {"type":"SESSION_END"}
  C->>S: F
  C->>S: reconnect fresh socket
  S->>S: reset_for_new_session(queue/event only)
""",
    ),
    SlideSpec(
        title="TCP 프로토콜",
        bullets=[
            "App에서 Server로 가는 신호는 P, S, X, E, N, F 여섯 가지다.",
            "Server에서 App으로 가는 제어 메시지는 START와 SESSION_END 두 가지 JSON이다.",
            "메시지는 newline-delimited JSON 계약을 사용한다.",
            "signal queue는 연속 XML만 latest로 collapse하고 제어 신호는 보존한다.",
        ],
        notes=[
            "근거: §3 TCP 프로토콜.",
            "제어 신호를 드롭하지 않는 큐 의미가 복구와 종료 시맨틱을 지탱한다.",
        ],
        mermaid="""\
flowchart LR
  subgraph APP[Android App]
    P[P package ACK]
    S[S screenshot payload]
    X[X XML + activity + package]
    E[E external app]
    N[N no-change]
    F[F finish]
  end
  subgraph SRV[Python Server]
    START[START package]
    ENDMSG[SESSION_END]
  end
  APP -->|newline-delimited JSON / payload| SRV
  SRV -->|START / SESSION_END| APP
""",
    ),
    SlideSpec(
        title="수집 루프 전체 흐름",
        bullets=[
            "AccessibilityEvent 이후 App이 안정화와 외부 앱 여부를 먼저 판단한다.",
            "변화가 있으면 screenshot·XML·metadata를 보내고 Server가 최신 신호를 소비한다.",
            "ScreenMatcher가 page_key와 observation을 정한 뒤 Explorer가 다음 action을 고른다.",
            "실행 결과는 events.jsonl의 frame_index, page_key, observation_num으로 누적된다.",
        ],
        notes=[
            "근거: §3 수집 루프.",
            "아래 다이어그램은 사용자가 고정한 CORE Diagram A 그대로다.",
        ],
        mermaid=DIAGRAM_A,
    ),
    SlideSpec(
        title="ScreenMatcher: BM25 unique-page matching",
        bullets=[
            "구조 지문 pre-filter가 exact revisit를 먼저 short-circuit한다.",
            "interactable button·input이 0개면 pending으로 거부해 빈 page blackhole을 막는다.",
            "새 화면은 element-line 문서로 직렬화해 BM25 top-K 후보를 검색한다.",
            "element 기준과 pixel gate를 함께 통과한 첫 후보를 기존 page로 병합한다.",
        ],
        notes=[
            "근거: §2 screen_matching, §3 수집 루프.",
            "아래 다이어그램은 사용자가 고정한 CORE Diagram B 그대로다.",
        ],
        mermaid=DIAGRAM_B,
    ),
    SlideSpec(
        title="Observation 식별과 저장 판단",
        bullets=[
            "page 식별이 끝나면 같은 page 안에서만 luminance 지문으로 observation을 다시 찾는다.",
            "같은 observation이면 재사용하고, 없으면 새 observation_num을 할당한다.",
            "is_new_observation일 때만 data/{package}/pages/{page_key}/{observation_num}에 새 파일을 쓴다.",
            "persist_filtered 기본값이 on이라 필터된 재방문도 방문마다 새 observation으로 저장된다.",
        ],
        notes=[
            "근거: §2 screen_matching, §5 저장 포맷.",
            "page 정체성과 observation 저장 판단이 분리돼 있다는 점을 강조한다.",
        ],
        mermaid="""\
flowchart TD
  PK[page_key fixed] --> LOOKUP[page-local luminance lookup]
  LOOKUP --> HIT{pixel diff < screenshot threshold?}
  HIT -->|yes| REUSE[reuse observation_num]
  HIT -->|no| NEWOBS[allocate next observation_num]
  REUSE --> WRITE0{is_new_observation?}
  NEWOBS --> WRITE1{is_new_observation?}
  WRITE0 -->|no| NOWRITE[skip save_observation]
  WRITE1 -->|yes| SAVE[write screenshot/XML/elements]
  PF[persist_filtered = on] --> NEWOBS
""",
    ),
    SlideSpec(
        title="LLMGuidedExplorer: select_action 4단계",
        bullets=[
            "진행 중 navigation이 있으면 queued step을 현재 화면에서 다시 매칭해 실행한다.",
            "현재 화면에 미탐색 action이 있으면 그중 하나를 고르고 long_touch는 후순위다.",
            "없으면 전역 미탐색 target을 골라 shortest path의 첫 step을 실행한다.",
            "마지막까지 없으면 back으로 후퇴하고 루트 화면에서는 in-screen tap을 쓴다.",
        ],
        notes=[
            "근거: §3 Action Space 와 탐색 전략.",
            "아래 다이어그램은 사용자가 고정한 CORE Diagram C 그대로다.",
        ],
        mermaid=DIAGRAM_C,
    ),
    SlideSpec(
        title="navigate-target 전략과 Action Space",
        bullets=[
            "전역 target 선택은 BFS, DFS, GREEDY 중 strategy 값으로 달라진다.",
            "BFS는 root에 가까운 화면, DFS는 더 깊은 화면, GREEDY는 가장 짧은 path를 우선한다.",
            "엔진 semantic action은 touch, select, long_touch, set_text, scroll이다.",
            "domain action은 Tap, LongPress, InputText, Swipe, PressBack이며 open_app은 record-only다.",
        ],
        notes=[
            "근거: §3 navigate-target 선택 전략, Action Space.",
            "기본 strategy는 config 기준 BFS이고, open_app은 navigation 전이에 쓰이지 않는다.",
        ],
        mermaid="""\
flowchart LR
  TARGET[global unexplored target] --> STRAT{strategy}
  STRAT --> BFS[BFS shallow target]
  STRAT --> DFS[DFS deep target]
  STRAT --> GREEDY[GREEDY shortest path]
  BFS --> PATH[TransitionGraph shortest path]
  DFS --> PATH
  GREEDY --> PATH
  PATH --> SEM[semantic action\n touch/select/long_touch/set_text/scroll]
  SEM --> DOM[domain action\n Tap/LongPress/InputText/Swipe/PressBack]
  REC[open_app record-only] -. transition false .-> DOM
""",
    ),
    SlideSpec(
        title="LLM 사용처",
        bullets=[
            "공용 LLMClient는 OpenRouter Chat Completions와 기본 모델 qwen/qwen3.7-plus를 쓴다.",
            "기본 on 사용처는 input-text 생성이고, 앱 설명을 프롬프트에 함께 넣는다.",
            "element 추출은 opt-in 기본 off이며 name·description·parameters·indexes를 한 번에 뽑는다.",
            "BM25 기반 page 식별 자체는 LLM-free라서 key가 없어도 dedup 흐름은 유지된다.",
        ],
        notes=[
            "근거: §2 llm, §2 screen_matching.",
            "LLM은 입력 생성과 element enrichment 두 곳만 쓰고 page identity는 비LLM이라는 점이 핵심이다.",
        ],
        mermaid="""\
flowchart TD
  CLIENT[LLMClient\nOpenRouter Chat Completions\nqwen/qwen3.7-plus] --> TXT[input-text generation]
  CLIENT --> EXT[element extraction opt-in]
  TXT --> PROMPT[App under test context]
  EXT --> ELEM[name/description/parameters\n element_index/key_element_index]
  BM25[BM25 unique-page matching] --> PAGE[page_key]
  BM25 -. LLM-free .-> PAGE
""",
    ),
    SlideSpec(
        title="저장 구조",
        bullets=[
            "data/{package}는 pages와 page_graph.json·page_graph.html 같은 영속 코퍼스를 담는다.",
            "runtime/{package}는 metadata, events.jsonl, coverage, cost 같은 휘발성 상태를 담는다.",
            "page_key와 observation_num은 0-based 정수 디렉터리명이며 조인 키로 재사용된다.",
            "converter와 오프라인 재구성은 events.jsonl의 page_key·observation_num을 기준으로 화면을 찾는다.",
        ],
        notes=[
            "근거: §5 저장 포맷.",
            "영속 root와 runtime root를 분리해 resume와 후처리 경로를 유지한다.",
        ],
        mermaid="""\
flowchart LR
  ROOT[data package root] --> PAGES[pages/page_key/observation_num]
  ROOT --> GRAPH[page_graph.json + page_graph.html]
  RUN[runtime package root] --> META[metadata.json]
  RUN --> EVENTS[events.jsonl]
  RUN --> COV[activity_coverage.csv]
  RUN --> COST[cost.csv]
  EVENTS --> JOIN[page_key + observation_num join]
  JOIN --> PAGES
""",
    ),
    SlideSpec(
        title="세션 관리와 복구",
        bullets=[
            "기본 동작은 resume이며 rehydrate_session이 page_graph와 ScreenMatcher 지식을 복원한다.",
            "MAX_NO_CHANGE_RETRIES=3, MAX_EXTERNAL_APP_RETRIES=10, MAX_SAME_PAGE_STEPS=5, MAX_EMPTY_UI_RETRIES=2를 쓴다.",
            "권한 다이얼로그는 grant 우선으로 자동 처리하고, external app 이탈은 return_to_app과 recover로 복구한다.",
            "타깃 앱을 다시 띄운 경우 open_app을 events.jsonl에 excursion당 1회 기록하되 transition:false로 격리한다.",
        ],
        notes=[
            "근거: §4 세션 관리와 복구.",
            "resume과 recovery는 데이터 손실을 막고 가짜 navigation 전이를 차단하는 방향으로 설계돼 있다.",
        ],
        mermaid="""\
flowchart TD
  START[resume session] --> REHYDRATE[rehydrate page_graph + ScreenMatcher]
  REHYDRATE --> LOOP[collection loop]
  LOOP --> NOCHANGE[no-change retry max 3]
  LOOP --> EMPTY[empty UI retry max 2]
  LOOP --> SAME[same-page guard max 5]
  LOOP --> EXT[external app retry max 10]
  EXT --> RETURN[return_to_app or recover]
  RETURN --> OPEN[log open_app once\ntransition false]
  PERM[permission dialog] --> GRANT[auto-grant preferred button]
""",
    ),
    SlideSpec(
        title="설정 시스템",
        bullets=[
            "설정은 builtin defaults, run.yaml, MC_* env, CLI 플래그 순으로 해석되고 later wins다.",
            "collection.budget_mode 기본값은 time, max_duration 기본값은 2h다.",
            "exploration.strategy의 canonical 기본값은 BFS이고 production 기본값도 BFS다.",
            "screen_matching 기본값은 luminance_prefilter=true, persist_filtered=true, bm25_top_k=5다.",
        ],
        notes=[
            "근거: §7 설정 시스템, config/run.yaml.",
            "문서상 기본값과 설정 우선순위를 한 장으로 정리한다.",
        ],
        mermaid="""\
flowchart LR
  B[builtin defaults] --> Y[config/run.yaml]
  Y --> E[MC_* env vars]
  E --> C[CLI flags]
  C --> CFG[resolved RunConfig]
  CFG --> BUDGET[budget_mode = time]
  CFG --> STRAT[strategy = BFS]
  CFG --> MATCH[luminance_prefilter = true\npersist_filtered = true\nbm25_top_k = 5]
""",
    ),
]


def render_mermaid_png(source: str) -> bytes:
    """Render mermaid to a PNG via kroki. Hard-fail (raise) on any failure —
    no local fallback renderer, so a broken render cannot silently masquerade
    as a real diagram."""
    request = urllib.request.Request(
        KROKI_URL,
        data=source.encode("utf-8"),
        headers={
            "Content-Type": "text/plain; charset=utf-8",
            "User-Agent": USER_AGENT,
        },
        method="POST",
    )
    last_error: Exception | None = None
    for _ in range(3):
        try:
            with urllib.request.urlopen(request, timeout=30) as response:
                data = response.read()
            if not data.startswith(b"\x89PNG\r\n\x1a\n"):
                raise ValueError("kroki response was not a PNG")
            return data
        except (urllib.error.URLError, TimeoutError, ValueError) as exc:
            last_error = exc
    raise RuntimeError(f"kroki render failed: {last_error}")


def fit_picture(slide, png_bytes: bytes) -> None:
    with Image.open(io.BytesIO(png_bytes)) as img:
        px_w, px_h = img.size
    area_left = Inches(6.2)
    area_top = Inches(1.55)
    area_w = Inches(6.6)
    area_h = Inches(5.45)
    ratio = min(area_w / px_w, area_h / px_h)
    width = int(px_w * ratio)
    height = int(px_h * ratio)
    left = int(area_left + (area_w - width) / 2)
    top = int(area_top + (area_h - height) / 2)
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        tmp.write(png_bytes)
        tmp_path = tmp.name
    slide.shapes.add_picture(tmp_path, left, top, width=width, height=height)


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    box.text_frame.text = text


def add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(5.1), Inches(5.7))
    box.text_frame.text = "\n".join(f"• {bullet}" for bullet in bullets)


def add_cover(slide, spec: SlideSpec) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(11.0), Inches(1.0))
    title_box.text_frame.text = "Monkey-Collector 수집 방식\n현재 코드·세팅 기준 파이프라인"
    body = slide.shapes.add_textbox(Inches(0.95), Inches(3.1), Inches(10.5), Inches(2.2))
    body.text_frame.text = "\n".join(f"• {bullet}" for bullet in spec.bullets)


def add_notes(slide, spec: SlideSpec) -> None:
    tf = slide.notes_slide.notes_text_frame
    text = "\n".join(spec.notes)
    if spec.mermaid:
        text = f"{text}\nMermaid source:\n{spec.mermaid.strip()}"
    tf.text = text


def build() -> Path:
    ROOT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    for index, spec in enumerate(SLIDES):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if index == 0:
            add_cover(slide, spec)
            add_notes(slide, spec)
            continue
        add_title(slide, spec.title)
        add_bullets(slide, spec.bullets)
        if spec.mermaid is None:
            raise ValueError(f"diagram slide {index} is missing mermaid")
        png = render_mermaid_png(spec.mermaid)
        fit_picture(slide, png)
        add_notes(slide, spec)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    output = build()
    print(output)
